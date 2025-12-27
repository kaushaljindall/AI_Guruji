import os
import uuid

# Safety Wrapper for SlideService
HAS_PPTX = False
HAS_PIL = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    print("⚠️ python-pptx not found. PPTX generation will be skipped.")

try:
    from PIL import Image, ImageDraw, ImageFont
    HAS_PIL = True
except ImportError:
    print("⚠️ Pillow not found. Slide Image generation will fail.")

class SlideService:
    def __init__(self):
        self.output_dir = os.path.join(os.getcwd(), "data", "outputs", "slides")
        os.makedirs(self.output_dir, exist_ok=True)

    def generate_presentation(self, lecture_title: str, slides_data: list[dict]) -> str:
        if not HAS_PPTX:
             print("❌ python-pptx missing. Skipping PPTX generation.")
             return ""
        
        try:
            prs = Presentation()
            
            # 1. Title Slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = lecture_title
            subtitle.text = "AI Guruji Lecture | By Kaushal Jindal"

            # 2. Content Slides
            bullet_slide_layout = prs.slide_layouts[1] 
            
            for index, slide_info in enumerate(slides_data):
                try:
                    slide = prs.slides.add_slide(bullet_slide_layout)
                    shapes = slide.shapes
                    title_shape = shapes.title
                    
                    # Fix: layout 1 content placeholder is usually index 1
                    body_shape = shapes.placeholders[1]
                    
                    # Title
                    title_shape.text = slide_info.get("heading", f"Slide {index+1}")
                    
                    # Content (Summary + Points)
                    tf = body_shape.text_frame
                    tf.clear() 
                    
                    summary = slide_info.get("summary", "")
                    if summary:
                        p = tf.add_paragraph()
                        p.text = summary
                        p.font.bold = True
                        p.space_after = Pt(14)
                    
                    points = slide_info.get("important_points", [])
                    if isinstance(points, str):
                        points = [points]
                        
                    for point in points:
                        p = tf.add_paragraph()
                        p.text = f"• {point}"
                        p.level = 0
                        p.space_after = Pt(10)
                        
                    # Code: Add as separate text box
                    code = slide_info.get("code", "")
                    if code:
                        left = Inches(0.5)
                        top = Inches(5.0)
                        width = Inches(9.0)
                        height = Inches(2.0)
                        txBox = slide.shapes.add_textbox(left, top, width, height)
                        tf_code = txBox.text_frame
                        p_code = tf_code.add_paragraph()
                        p_code.text = code
                        p_code.font.name = "Courier New"
                        p_code.font.size = Pt(11)

                except Exception as e:
                    print(f"Error creating slide {index}: {e}")
                    continue 

            filename = f"{lecture_title.replace(' ', '_')}_{uuid.uuid4()}.pptx"
            filepath = os.path.join(self.output_dir, filename)
            prs.save(filepath)
            print(f"Presentation saved to {filepath}")
            return filepath
            
        except Exception as e:
            print(f"Critical Error in generate_presentation: {e}")
            return ""

    def generate_slide_image(self, slide_data: dict, index: int) -> str:
        """
        Generates an image of the slide using Pillow.
        If Pillow is missing, returns empty string.
        """
        if not HAS_PIL:
            print("❌ Pillow missing. Cannot generate slide image.")
            return ""

        try:
            width, height = 1920, 1080
            img = Image.new('RGB', (width, height), color='#1e1e1e')
            draw = ImageDraw.Draw(img)
            
            # Fonts
            try:
                title_font = ImageFont.truetype("arialbd.ttf", 70)
                body_font = ImageFont.truetype("arial.ttf", 40)
                small_font = ImageFont.truetype("consola.ttf", 30)
            except IOError:
                try:
                     title_font = ImageFont.truetype("DejaVuSans-Bold.ttf", 70)
                     body_font = ImageFont.truetype("DejaVuSans.ttf", 40)
                     small_font = ImageFont.truetype("DejaVuSansMono.ttf", 30)
                except:
                     title_font = ImageFont.load_default()
                     body_font = ImageFont.load_default()
                     small_font = ImageFont.load_default()

            margin_x = 100
            current_y = 100
            
            # 1. Heading
            heading = slide_data.get("heading", f"Slide {index+1}")
            draw.text((margin_x, current_y), heading, font=title_font, fill="#60a5fa")
            current_y += 120
            
            draw.line((margin_x, current_y, width - margin_x, current_y), fill="#334155", width=4)
            current_y += 60
            
            # 2. Summary
            summary = slide_data.get("summary", "")
            if summary:
                self._draw_text_wrapped(draw, summary, 20, body_font, margin_x, current_y, width - margin_x)
                current_y += 150 
            
            # 3. Points
            points = slide_data.get("important_points", [])
            if isinstance(points, str): points = [points]
                
            for point in points:
                text = f"• {point}"
                draw.text((margin_x, current_y), text, font=body_font, fill="#e2e8f0")
                current_y += 60
                
            # 4. Code
            code = slide_data.get("code", "")
            if code:
                current_y += 40
                draw.rectangle((margin_x, current_y, width - margin_x, height - 100), fill="#0f172a", outline="#334155")
                draw.text((margin_x + 20, current_y + 20), code, font=small_font, fill="#d1d5db")

            filename = f"slide_{index}_{uuid.uuid4()}.png"
            filepath = os.path.join(self.output_dir, filename)
            img.save(filepath)
            return filepath
            
        except Exception as e:
            print(f"Error drawing slide image: {e}")
            return ""

    def _draw_text_wrapped(self, draw, text, char_width, font, x, y, max_width):
        # ... logic consistent with previous ...
        lines = []
        words = text.split()
        current_line = []
        for word in words:
            current_line.append(word)
            if len(" ".join(current_line)) * char_width * 0.5 > max_width: 
                current_line.pop()
                lines.append(" ".join(current_line))
                current_line = [word]
        lines.append(" ".join(current_line))
        
        for line in lines:
            draw.text((x, y), line, font=font, fill="#ffffff")
            y += 50
        return y

slide_service = SlideService()
